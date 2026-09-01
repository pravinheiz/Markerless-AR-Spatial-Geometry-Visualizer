import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    C_BG = RGBColor(11, 17, 30)        # #0B111E Dark Tech Slate
    C_CARD = RGBColor(18, 28, 48)      # #121C30 Glass Card
    C_CARD_BORDER = RGBColor(30, 45, 75)
    C_PRIMARY = RGBColor(0, 255, 136)  # #00FF88 Neon Green
    C_ACCENT = RGBColor(0, 204, 255)   # #00CCFF Cyan
    C_WHITE = RGBColor(255, 255, 255)
    C_MUTED = RGBColor(160, 180, 205)
    C_WARN = RGBColor(255, 170, 0)     # #FFAA00 Gold

    def add_header(slide, tagline, title, highlight_word=""):
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()

        # Top Accent Line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = C_PRIMARY
        top_line.line.fill.background()

        # Tagline
        tx_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.4))
        p_tag = tx_tag.text_frame.paragraphs[0]
        p_tag.text = tagline.upper()
        p_tag.font.name = 'Arial'
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C_ACCENT

        # Title
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.7))
        p_title = tx_title.text_frame.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = C_WHITE

    # SLIDE 1: Title & Overview
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_BG
    bg1.line.fill.background()

    # Title Hero Box
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.2))
    tf = t_box.text_frame
    p0 = tf.paragraphs[0]
    p0.text = "OPEN PROJECT LAB DEMONSTRATION // ASSIGNMENT SUBMISSION"
    p0.font.name = 'Arial'
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = C_ACCENT

    p1 = tf.add_paragraph()
    p1.text = "Markerless AR Spatial Geometry Visualizer"
    p1.font.name = 'Arial'
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = C_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "Next-Generation Browser-Based Touchless 3D Holographic Modeling & Analytics"
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_MUTED

    # 3 Summary Cards
    cards = [
        ("🖐️ Markerless Computer Vision", "Real-time 21-keypoint 3D hand tracking powered by client-side MediaPipe ML neural network (60 FPS, zero hardware)."),
        ("✨ 6-DOF Tony Stark Control", "True mid-air freeform 3D wrist-tilt orientation (Pitch, Yaw, Roll) & dual-hand Iron Man pinch zoom scaling."),
        ("📊 Live Spatial Telemetry HUD", "Real-time geometric computation of Volume (V), Surface Area (A), Euler topology, and 2-hand laser ruler distance.")
    ]
    for i, (head, desc) in enumerate(cards):
        c_shape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0 + i * 3.9), Inches(3.8), Inches(3.6), Inches(2.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        c_shape.line.width = Pt(1.5)
        
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = C_WHITE
        
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = C_MUTED

    s1.notes_slide.notes_text_frame.text = "Slide 1: Welcome to the presentation of Markerless AR Spatial Geometry Visualizer. This project brings touchless, gesture-driven 3D CAD and spatial geometry analysis directly into modern web browsers."

    # SLIDE 2: Problem Statement
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "The Problem & Academic Motivation", "Overcoming the 2D Flatland Barrier in Spatial Computing")
    p_cards = [
        ("🖥️ 2D Screen Constraint", "Traditional CAD and geometry visualization tools force 3D volumetric ideas onto flat 2D monitors, creating significant cognitive friction in spatial comprehension."),
        ("🥽 Costly AR/VR Hardware", "Dedicated headsets (Apple Vision Pro, Meta Quest 3, HoloLens) cost $500–$3,500, creating severe accessibility barriers for schools, students, and researchers."),
        ("⚙️ Fiducial Marker Friction", "Conventional marker-based AR depends on printed paper QR codes or ArUco tags, restricting natural hand mobility and freeform mid-air interaction.")
    ]
    for i, (head, desc) in enumerate(p_cards):
        c_shape = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_PRIMARY
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    s2.notes_slide.notes_text_frame.text = "Slide 2: Traditional spatial modeling is constrained by 2D screens, while VR headsets are too costly. Our project removes both obstacles."

    # SLIDE 3: Proposed Solution
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Innovative Solution", "The Markerless Web-AR Architecture")
    s_cards = [
        ("🌐 Universal Web Accessibility", "Runs natively in any modern browser via Three.js (WebGL) and MediaPipe (WebAssembly) with zero software installation, plugins, or driver setup."),
        ("⚡ Sub-20ms Real-Time Vision", "Client-side neural network tracking processes 21 3D hand keypoints per hand at a smooth 60 FPS, ensuring zero perceivable input latency."),
        ("🎯 Anti-Jitter Hysteresis Filter", "Temporal exponential smoothing filter eliminates hand tremor and prevents accidental drawing triggers when waving hands in front of the camera.")
    ]
    for i, (head, desc) in enumerate(s_cards):
        c_shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_ACCENT
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 4: Technical Architecture
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Technical Architecture", "Three-Tiered Synchronous Render Pipeline")
    layers = [
        ("LAYER 2: Holographic HUD Canvas", "Renders glowing skeletal joints, Oculus Quest index fingertip cursor, and 2-hand laser measurement telemetry overlays."),
        ("LAYER 1: Three.js 3D WebGL Canvas", "Transparent 3D scene managing lighting, shadow mapping, Catmull-Rom tube ribbons, shader materials, and raycasting."),
        ("LAYER 0: Webcam Video Stream", "Mirrored scaleX(-1) background stream with calibrated coordinate mapping to Three.js Normalized Device Coordinates.")
    ]
    for i, (head, desc) in enumerate(layers):
        c_shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8 + i * 1.65), Inches(11.733), Inches(1.4))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = C_PRIMARY if i==1 else C_ACCENT
        p_d = tf_c.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 5: Gesture Interaction Engine
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Natural User Interaction (NUI)", "Tony Stark 6-DOF & Oculus Quest Gesture Controls")
    g_cards = [
        ("✊ 6-DOF Wrist Orientation", "Pitch (wrist-to-finger tilt), Yaw (lateral angle), and Roll (thumb-to-pinky twist) rotate active 3D holograms in mid-air without predefined buttons."),
        ("🤏🤏 Iron Man Dual Zoom", "Pinch both hands and spread apart to expand/zoom holograms up to 5.0x; bring hands together to shrink/zoom out seamlessly."),
        ("🖐️ Oculus UI Touch", "Index fingertip projects an optical cursor. Hovering over UI buttons highlights them; pinching clicks buttons with synthesized audio feedback.")
    ]
    for i, (head, desc) in enumerate(g_cards):
        c_shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_PRIMARY
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 6: 3D AR Painting Engine
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Creative Spatial Tools", "Oculus / Samsung AR 3D Painting & Extrusion")
    paint_cards = [
        ("🎨 Volumetric Tube Ribbons", "Pinching in 3D Paint mode generates continuous, smooth Catmull-Rom tube geometries in spatial depth rather than flat 2D lines."),
        ("🌈 6-Color Shader Palette", "Paint with Neon Green, Cyan, Electric Purple, Gold, Hot Pink, and Pure White across Hologram, Glass, Metal, and Neon shaders."),
        ("✨ 2-Hand Spline Extrusion", "Two-hand vertical separation gesture automatically inflates 2D hand-sketched profiles into solid 3D spatial tube meshes.")
    ]
    for i, (head, desc) in enumerate(paint_cards):
        c_shape = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_ACCENT
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 7: Spatial Metrics Analytics
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Mathematical Computation", "Live Spatial Geometry Analytics HUD")
    m_rows = [
        ("Volume (V)", "V = ∭ dV (Exact for Platonic solids, cylinders, cones, toruses, and mesh tetrahedrons)", "Volumetric capacity and material weight estimation"),
        ("Surface Area (A)", "A = ∬ dA (Sum of all indexed triangular face areas)", "Surface coating, material cost, and heat dissipation"),
        ("Bounding Box", "ΔX × ΔY × ΔZ via THREE.BoxHelper dynamic projection", "Clearance checking, spatial envelope, and packaging sizing"),
        ("Euler Topology", "V - E + F = 2 (Vertices, Edges, Faces counter)", "Topological integrity analysis and mesh complexity evaluation")
    ]
    for i, (met, form, purp) in enumerate(m_rows):
        c_shape = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8 + i * 1.25), Inches(11.733), Inches(1.1))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = f"{met}: {form}"
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(14)
        p_h.font.bold = True
        p_h.font.color.rgb = C_PRIMARY
        p_d = tf_c.add_paragraph()
        p_d.text = "Purpose: " + purp
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 8: Measurement & Laser Ruler
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Measurement & Telemetry", "AR Dual-Hand Laser Distance Ruler")
    r_cards = [
        ("📏 3D Laser Dimension Line", "Projects a glowing AR dashed laser beam connecting the index fingertips of both hands in 3D spatial space."),
        ("📟 Real-Time Telemetry Pill", "Floating HUD overlay displays instant Euclidean metric distance readout in meters (e.g. DIST: 0.48 m)."),
        ("🎯 Sub-Centimeter Accuracy", "Transforms normalized camera depth into calibrated physical spatial units with high precision.")
    ]
    for i, (head, desc) in enumerate(r_cards):
        c_shape = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_PRIMARY
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 9: Applications & Benefits
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Impact & Benefits", "Real-World Industry & Academic Use Cases")
    u_cards = [
        ("🎓 STEM & Math Education", "Transforms abstract 3D calculus, geometry, and vector physics into tangible, interactive mid-air learning objects for classrooms."),
        ("🏭 Rapid CAD Prototyping", "Allows industrial designers and engineers to conceptualize volumetric 3D concepts in seconds before exporting to CAD suites."),
        ("🩺 Sterile Medical Visualization", "Enables surgeons and medical students to manipulate 3D anatomical organ models and CT scans in touchless sterile environments.")
    ]
    for i, (head, desc) in enumerate(u_cards):
        c_shape = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_ACCENT
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    # SLIDE 10: Future Roadmap & Conclusion
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Summary & Future Roadmap", "The Future of Touchless Web-Based Spatial AR")
    f_cards = [
        ("☁️ GLTF / OBJ Model Export", "One-click 3D model download pipeline allowing hand-drawn spatial geometry to be exported directly to 3D printing slicers."),
        ("👥 Multi-User WebRTC AR", "Real-time multi-peer collaboration allowing distributed teams to co-design in the same holographic space."),
        ("🚀 WebXR Headset Integration", "Seamless cross-compatibility bridging desktop webcams to Meta Quest 3, Apple Vision Pro, and mobile AR.")
    ]
    for i, (head, desc) in enumerate(f_cards):
        c_shape = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = C_CARD
        c_shape.line.color.rgb = C_CARD_BORDER
        tf_c = c_shape.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = C_PRIMARY
        p_d = tf_c.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_MUTED

    output_path = "e:/Assignments/GD/Lab/Project/Markerless_AR_Spatial_Geometry_Visualizer.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_deck()
